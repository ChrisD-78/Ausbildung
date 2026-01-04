#!/usr/bin/env python3
"""
Erstellt eine PowerPoint-Präsentation aus allen Folien der index.html
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from bs4 import BeautifulSoup

def extract_slides_from_html():
    """Extrahiert alle Folien aus der index.html"""
    with open('index.html', 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides = []
    
    # Finde alle Slide-Divs
    slide_divs = soup.find_all('div', class_='slide')
    
    for slide_div in slide_divs:
        slide_data = {
            'title': '',
            'subtitle': '',
            'content': []
        }
        
        # Extrahiere Titel (h2)
        h2 = slide_div.find('h2')
        if h2:
            slide_data['title'] = h2.get_text(strip=True)
        
        # Extrahiere Untertitel (h3)
        h3 = slide_div.find('h3')
        if h3:
            slide_data['subtitle'] = h3.get_text(strip=True)
        
        # Extrahiere weiteren Inhalt
        # Entferne bereits extrahierte Elemente
        if h2:
            h2.decompose()
        if h3:
            h3.decompose()
        
        # Extrahiere Text-Inhalte
        content_text = slide_div.get_text(separator='\n', strip=True)
        slide_data['content'] = [line for line in content_text.split('\n') if line.strip()]
        
        slides.append(slide_data)
    
    return slides

def create_presentation():
    """Erstellt die PowerPoint-Präsentation"""
    prs = Presentation()
    
    # Setze Präsentationsformat (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Extrahiere Folien aus HTML
    slides_data = extract_slides_from_html()
    
    for slide_data in slides_data:
        # Erstelle neue Folie
        slide_layout = prs.slide_layouts[6]  # Leere Folie
        slide = prs.slides.add_slide(slide_layout)
        
        # Hintergrundfarbe
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea
        
        # Titel hinzufügen
        if slide_data['title']:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.word_wrap = True
            
            # Titel-Formatierung
            for paragraph in title_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(44)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Untertitel hinzufügen
        if slide_data['subtitle']:
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9)
            height = Inches(0.8)
            
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data['subtitle']
            subtitle_frame.word_wrap = True
            
            for paragraph in subtitle_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(28)
                    run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Inhalt hinzufügen
        if slide_data['content']:
            left = Inches(0.5)
            top = Inches(3)
            width = Inches(9)
            height = Inches(4)
            
            content_box = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Füge Inhalt hinzu
            content_text = '\n'.join(slide_data['content'][:15])  # Maximal 15 Zeilen
            content_frame.text = content_text
            
            for paragraph in content_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Speichere Präsentation
    output_file = 'Südpfalztreffen Bäder.pptx'
    prs.save(output_file)
    print(f"✅ PowerPoint-Präsentation erstellt: {output_file}")
    return output_file

if __name__ == '__main__':
    try:
        create_presentation()
    except Exception as e:
        print(f"Fehler: {e}")
        import traceback
        traceback.print_exc()


